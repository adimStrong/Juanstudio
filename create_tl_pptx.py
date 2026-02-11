#!/usr/bin/env python3
"""
Create a TL-friendly PowerPoint presentation for the 15-Day Sprint Program.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== COLOR PALETTE =====
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT_YELLOW = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)
CARD_BORDER = RGBColor(0x33, 0x40, 0x55)


def add_dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_shape_with_fill(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(4)
    return txBox


def add_metric_card(slide, left, top, width, height, title, value, subtitle="", accent_color=ACCENT_BLUE):
    add_shape_with_fill(slide, left, top, width, height, CARD_BG, CARD_BORDER)
    add_rect(slide, left, top, width, Pt(4), accent_color)
    add_text_box(slide, left + Inches(0.2), top + Pt(12), width - Inches(0.4), Inches(0.3),
                 title, 11, LIGHT_GRAY, False)
    add_text_box(slide, left + Inches(0.2), top + Pt(36), width - Inches(0.4), Inches(0.4),
                 value, 24, WHITE, True)
    if subtitle:
        add_text_box(slide, left + Inches(0.2), top + Pt(72), width - Inches(0.4), Inches(0.3),
                     subtitle, 10, LIGHT_GRAY, False)


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_rect(slide, Inches(1.5), Inches(2.5), Inches(1), Pt(4), ACCENT_BLUE)
add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.2),
             "15-Day Sprint Program", 48, WHITE, True)
add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8),
             'Client Base ("Pond") Cultivation & Team Leader Selection', 24, LIGHT_GRAY)
add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.5),
             "Team Leader Briefing  |  Skyline Nexus Corporation  |  February 2026", 14, ACCENT_BLUE)


# ============================================================
# SLIDE 2: Program Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "PROGRAM OVERVIEW", 12, ACCENT_BLUE, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
             "What Is This Program?", 32, WHITE, True)

cards_data = [
    ("MISSION", "Transform 10 new agents into independent VIP Client Handlers with self-sustaining client bases within 15 days", ACCENT_BLUE),
    ("SELECTION", "Identify top 4-6 performers for immediate promotion to Team Leader", ACCENT_GREEN),
    ("SCALE", "10 agents -> 30 -> 90 agents in 45 days through repeating 15-day cycles", ACCENT_PURPLE),
]

for i, (title, desc, color) in enumerate(cards_data):
    left = Inches(0.8 + i * 4.0)
    top = Inches(1.8)
    add_shape_with_fill(slide, left, top, Inches(3.6), Inches(2.0), CARD_BG, CARD_BORDER)
    add_rect(slide, left, top, Inches(3.6), Pt(4), color)
    add_text_box(slide, left + Inches(0.2), top + Pt(14), Inches(3.2), Inches(0.3),
                 title, 12, color, True)
    add_text_box(slide, left + Inches(0.2), top + Pt(40), Inches(3.2), Inches(1.2),
                 desc, 14, LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(4.2), Inches(6), Inches(0.5),
             "Team Composition", 20, WHITE, True)

comp_items = [
    ("10 Agents", "5 Male + 5 Female", ACCENT_BLUE),
    ("30-40% Attrition", "Survival of the Fittest", ACCENT_RED),
    ("4-6 TL Promotions", "Top Performers Promoted", ACCENT_GREEN),
    ("15 Calendar Days", "3 Phases of Execution", ACCENT_YELLOW),
]

for i, (title, sub, color) in enumerate(comp_items):
    left = Inches(0.8 + i * 3.1)
    add_metric_card(slide, left, Inches(4.8), Inches(2.8), Inches(1.4), sub, title, "", color)


# ============================================================
# SLIDE 3: Day 15 Success Targets
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "DAY 15 TARGETS", 12, ACCENT_BLUE, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
             "Non-Negotiable Success Criteria", 32, WHITE, True)

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(6), Inches(0.4),
             "Per Agent Targets", 18, WHITE, True)

targets = [
    ("Active Pond Size", ">=50", "Clients", "Interacted in last 7 days", ACCENT_GREEN),
    ("Daily Net Growth", "+5 to +8", "Clients/Day", "New adds minus freezes", ACCENT_BLUE),
    ("Total Pond Capacity", "80-120", "Clients", "Active + Dormant (cap 200)", ACCENT_PURPLE),
    ("Independence Score", "100%", "Independent", "Zero escalations needed", ACCENT_YELLOW),
]

for i, (title, value, unit, desc, color) in enumerate(targets):
    left = Inches(0.8 + i * 3.1)
    add_shape_with_fill(slide, left, Inches(2.2), Inches(2.8), Inches(2.0), CARD_BG, CARD_BORDER)
    add_rect(slide, left, Inches(2.2), Inches(2.8), Pt(4), color)
    add_text_box(slide, left + Inches(0.15), Inches(2.35), Inches(2.5), Inches(0.3), title, 12, LIGHT_GRAY)
    add_text_box(slide, left + Inches(0.15), Inches(2.7), Inches(2.5), Inches(0.5), value, 36, WHITE, True)
    add_text_box(slide, left + Inches(0.15), Inches(3.35), Inches(2.5), Inches(0.3), unit, 14, color, True)
    add_text_box(slide, left + Inches(0.15), Inches(3.65), Inches(2.5), Inches(0.3), desc, 10, LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(4.6), Inches(6), Inches(0.4),
             "Failure = Immediate Re-assignment or Termination", 18, ACCENT_RED, True)

add_shape_with_fill(slide, Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.5), CARD_BG, ACCENT_RED)
add_multi_text(slide, Inches(1.1), Inches(5.25), Inches(11), Inches(1.2), [
    ("FAIL Condition 1:  Active Pond Size < 40 clients on Day 15", 14, WHITE, False),
    ("FAIL Condition 2:  Average net growth over last 5 days <= 0 (shrinking pond)", 14, WHITE, False),
    ("Result:  Immediate re-assignment to standard accounts or termination", 14, ACCENT_RED, True),
])


# ============================================================
# SLIDE 4: 3 Phases Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "15-DAY EXECUTION", 12, ACCENT_BLUE, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
             "Three Phases of the Sprint", 32, WHITE, True)

phases = [
    ("PHASE 1", "Days 1-3", "Foundation & Assimilation",
     "Learn the Tools, Master the Process",
     ["Day 1: System immersion, receive 50 clients, shadow 3 calls",
      "Day 2: Call all 50 clients, first replenishment",
      "Day 3: Focus responsive clients, coaching on tone"],
     "Target: 20-25 Active Clients", ACCENT_BLUE),
    ("PHASE 2", "Days 4-10", "Accelerated Growth",
     "Build Velocity, Optimize Value",
     ["Days 4-7: 60% maintenance calls, 40% acquisition",
      "Days 8-10: Identify Top 20 clients, tiered engagement",
      "Calls increase to 40-45/day"],
     "Target: 45-55 Active Clients", ACCENT_GREEN),
    ("PHASE 3", "Days 11-15", "Stabilization & Selection",
     "Quality & Leadership Display",
     ["Days 11-13: Rebalance pond, strategic adds only",
      "Days 14-15: Final sprint, leadership testing",
      "Quality over quantity focus"],
     "Target: 60+ Active Clients (TL)", ACCENT_PURPLE),
]

for i, (phase, days, title, theme, bullets, target, color) in enumerate(phases):
    left = Inches(0.5 + i * 4.2)
    add_shape_with_fill(slide, left, Inches(1.7), Inches(3.9), Inches(5.2), CARD_BG, CARD_BORDER)
    add_shape_with_fill(slide, left, Inches(1.7), Inches(3.9), Inches(1.0), color)
    add_text_box(slide, left + Inches(0.2), Inches(1.8), Inches(3.5), Inches(0.3),
                 f"{phase}  |  {days}", 12, WHITE, True)
    add_text_box(slide, left + Inches(0.2), Inches(2.1), Inches(3.5), Inches(0.4),
                 title, 18, WHITE, True)
    add_text_box(slide, left + Inches(0.2), Inches(2.85), Inches(3.5), Inches(0.3),
                 theme, 11, color, False)
    for j, bullet in enumerate(bullets):
        add_text_box(slide, left + Inches(0.2), Inches(3.3 + j * 0.55), Inches(3.5), Inches(0.5),
                     f"  {bullet}", 11, LIGHT_GRAY)
    add_shape_with_fill(slide, left + Inches(0.15), Inches(5.9), Inches(3.6), Inches(0.6), color)
    add_text_box(slide, left + Inches(0.3), Inches(5.95), Inches(3.3), Inches(0.5),
                 target, 14, WHITE, True, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5: Daily Targets Ramp
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "DAILY TARGETS", 12, ACCENT_BLUE, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
             "Day-by-Day Agent Targets", 32, WHITE, True)

headers = ["Day", "Calls", "New Adds", "Active Target", "Net Growth", "Phase"]
col_widths = [Inches(1.2), Inches(1.5), Inches(1.5), Inches(2.0), Inches(1.5), Inches(2.5)]
col_starts = [Inches(0.8)]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

add_shape_with_fill(slide, Inches(0.8), Inches(1.6), Inches(10.2), Inches(0.45), ACCENT_BLUE)
for i, (header, left, width) in enumerate(zip(headers, col_starts, col_widths)):
    add_text_box(slide, left, Inches(1.65), width, Inches(0.35), header, 12, WHITE, True)

rows = [
    ("Day 1", "20-25", "0", "~15", "+0", "Phase 1: Foundation", ACCENT_BLUE),
    ("Day 2", "30", "3-5", ">=15", "+3-5", "Phase 1: Foundation", ACCENT_BLUE),
    ("Day 3", "35", "5-7", ">=20", "+5", "Phase 1: Foundation", ACCENT_BLUE),
    ("Day 4-7", "40", "8-10", ">=35", "+5/day", "Phase 2: Growth", ACCENT_GREEN),
    ("Day 8-10", "45", "3-5 HV", ">=45", "+5/day", "Phase 2: Optimization", ACCENT_GREEN),
    ("Day 11-13", "40", "Strategic", ">=50", "+3-5", "Phase 3: Stabilize", ACCENT_PURPLE),
    ("Day 14-15", "35", "Quality", ">=60 (TL)", "+2-3", "Phase 3: Assessment", ACCENT_PURPLE),
]

for j, (day, calls, adds, active, growth, phase, color) in enumerate(rows):
    row_top = Inches(2.1 + j * 0.55)
    bg_color = CARD_BG if j % 2 == 0 else DARK_BG
    add_shape_with_fill(slide, Inches(0.8), row_top, Inches(10.2), Inches(0.5), bg_color)
    add_rect(slide, Inches(0.8), row_top, Pt(4), Inches(0.5), color)
    values = [day, calls, adds, active, growth, phase]
    for i, (val, left, width) in enumerate(zip(values, col_starts, col_widths)):
        txt_color = WHITE if i == 3 else LIGHT_GRAY
        txt_bold = True if i == 3 else False
        add_text_box(slide, left + Inches(0.1), row_top + Pt(4), width, Inches(0.4),
                     val, 12, txt_color, txt_bold)

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
             "Key: Active clients = interacted within last 7 days  |  Frozen = no interaction for 7+ days = removed from pond",
             12, ACCENT_YELLOW)


# ============================================================
# SLIDE 6: Iron Rules
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "IRON RULES", 12, ACCENT_RED, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
             "Non-Negotiable Rules & Consequences", 32, WHITE, True)

rules = [
    ("1", "7-Day Freeze Rule",
     "No interaction for 7 days = client FROZEN and removed from pond",
     "Automated daily flag at 9 AM. Manager approval for exceptions.",
     ACCENT_RED),
    ("2", "Daily Growth Mandate",
     "Must add enough new clients to replace freezes + 3-5 extra daily",
     "Fail 2 consecutive days = probation. 3 days = removal.",
     ACCENT_YELLOW),
    ("3", "100% Activity Logging",
     "Every call must be logged in CRM. Random audit of 20% daily.",
     "Compliance < 95% = ZERO commission for the day.",
     ACCENT_BLUE),
    ("4", "Zero-Tolerance Compliance",
     "Any client complaint triggers immediate investigation.",
     "Verified violation = immediate dismissal. No exceptions.",
     ACCENT_RED),
]

for i, (num, title, desc, consequence, color) in enumerate(rules):
    top = Inches(1.7 + i * 1.35)
    add_shape_with_fill(slide, Inches(0.8), top, Inches(11.7), Inches(1.2), CARD_BG, CARD_BORDER)
    num_shape = add_shape_with_fill(slide, Inches(1.0), top + Inches(0.25), Inches(0.6), Inches(0.6), color)
    add_text_box(slide, Inches(1.0), top + Inches(0.28), Inches(0.6), Inches(0.5),
                 num, 24, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.9), top + Inches(0.1), Inches(4), Inches(0.4),
                 title, 18, WHITE, True)
    add_text_box(slide, Inches(1.9), top + Inches(0.55), Inches(5), Inches(0.5),
                 desc, 12, LIGHT_GRAY)
    add_text_box(slide, Inches(7.2), top + Inches(0.1), Inches(5), Inches(0.4),
                 "CONSEQUENCE:", 10, color, True)
    add_text_box(slide, Inches(7.2), top + Inches(0.45), Inches(5), Inches(0.6),
                 consequence, 12, LIGHT_GRAY)


# ============================================================
# SLIDE 7: Daily Management
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "DAILY MANAGEMENT", 12, ACCENT_BLUE, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
             "TL Daily Management Protocol", 32, WHITE, True)

# Morning Huddle
add_shape_with_fill(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(3.5), CARD_BG, CARD_BORDER)
add_rect(slide, Inches(0.8), Inches(1.7), Inches(5.5), Pt(4), ACCENT_YELLOW)
add_text_box(slide, Inches(1.0), Inches(1.85), Inches(5), Inches(0.4),
             "Morning Huddle  (08:45 - 08:55)", 18, ACCENT_YELLOW, True)
add_multi_text(slide, Inches(1.0), Inches(2.35), Inches(5), Inches(2.5), [
    ("1. Review yesterday's Ranking Score:", 13, WHITE, True),
    ("   Score = (Active*0.4) + (Growth*0.3) + (Quality*0.2) + (Compliance*0.1)", 11, LIGHT_GRAY, False),
    ("2. Announce Top 3 & Bottom 3 agents", 13, WHITE, True),
    ("3. Set daily growth target:", 13, WHITE, True),
    ("   Team freezes yesterday + 40 = new client target", 11, LIGHT_GRAY, False),
])

# Monitoring Dashboard
add_shape_with_fill(slide, Inches(6.8), Inches(1.7), Inches(5.7), Inches(3.5), CARD_BG, CARD_BORDER)
add_rect(slide, Inches(6.8), Inches(1.7), Inches(5.7), Pt(4), ACCENT_GREEN)
add_text_box(slide, Inches(7.0), Inches(1.85), Inches(5), Inches(0.4),
             "Real-Time KPI Dashboard", 18, ACCENT_GREEN, True)

kpis = [
    ("Active Client Ratio", "Active / Total Pond", ">=60%", "<=40%"),
    ("Growth Velocity", "Slope of net growth", ">=0.5", "<=0"),
    ("Call Efficiency", "Adds / Calls Made", ">=1:8", "<=1:15"),
    ("Client Value Index", "Revenue / Active", "Rising", "Falling"),
]
for j, (name, formula, green, red) in enumerate(kpis):
    y = Inches(2.4 + j * 0.6)
    add_text_box(slide, Inches(7.0), y, Inches(2.5), Inches(0.3), name, 11, WHITE, True)
    add_text_box(slide, Inches(9.5), y, Inches(1.5), Inches(0.3), green, 11, ACCENT_GREEN, True)
    add_text_box(slide, Inches(11.0), y, Inches(1.5), Inches(0.3), red, 11, ACCENT_RED, True)

# End of Day
add_shape_with_fill(slide, Inches(0.8), Inches(5.5), Inches(5.5), Inches(1.6), CARD_BG, CARD_BORDER)
add_rect(slide, Inches(0.8), Inches(5.5), Inches(5.5), Pt(4), ACCENT_BLUE)
add_text_box(slide, Inches(1.0), Inches(5.65), Inches(5), Inches(0.4),
             "End-of-Day Agent Log (Due by 21:00)", 16, ACCENT_BLUE, True)
add_multi_text(slide, Inches(1.0), Inches(6.1), Inches(5), Inches(0.9), [
    ("Agents submit: Calls, Freezes, New Adds, Net Growth, Active Count", 11, LIGHT_GRAY, False),
    ("TL spot-checks 2 random logs against CRM daily", 11, WHITE, True),
])

# QA
add_shape_with_fill(slide, Inches(6.8), Inches(5.5), Inches(5.7), Inches(1.6), CARD_BG, CARD_BORDER)
add_rect(slide, Inches(6.8), Inches(5.5), Inches(5.7), Pt(4), ACCENT_PURPLE)
add_text_box(slide, Inches(7.0), Inches(5.65), Inches(5), Inches(0.4),
             "Quality Assurance (Daily)", 16, ACCENT_PURPLE, True)
add_multi_text(slide, Inches(7.0), Inches(6.1), Inches(5), Inches(0.9), [
    ("3 random call recordings reviewed per agent per day", 11, LIGHT_GRAY, False),
    ("Score < 7.5 = next-day coaching session", 11, WHITE, True),
    ("Rubric: Tone(30%) + Process(30%) + Value(20%) + Upsell(20%)", 11, LIGHT_GRAY, False),
])


# ============================================================
# SLIDE 8: TL Selection Matrix
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "TL SELECTION", 12, ACCENT_GREEN, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
             "Team Leader Selection Matrix", 32, WHITE, True)

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(6), Inches(0.4),
             "Quantitative Criteria (Must Meet ALL)", 18, ACCENT_BLUE, True)

quant = [
    ("Active Pond (D15)", ">=60 Clients", "25%", ACCENT_GREEN),
    ("Growth Consistency", "7/12 days >=+5", "20%", ACCENT_BLUE),
    ("Client Quality", "Top 30%", "20%", ACCENT_PURPLE),
    ("Compliance", ">=98%", "15%", ACCENT_YELLOW),
    ("QA Score", ">=8.5/10", "20%", ACCENT_RED),
]

for i, (name, target, weight, color) in enumerate(quant):
    left = Inches(0.8 + i * 2.4)
    add_shape_with_fill(slide, left, Inches(2.1), Inches(2.2), Inches(1.5), CARD_BG, CARD_BORDER)
    add_rect(slide, left, Inches(2.1), Inches(2.2), Pt(4), color)
    add_text_box(slide, left + Inches(0.1), Inches(2.2), Inches(2.0), Inches(0.3), name, 11, LIGHT_GRAY)
    add_text_box(slide, left + Inches(0.1), Inches(2.55), Inches(2.0), Inches(0.4), target, 18, WHITE, True)
    add_text_box(slide, left + Inches(0.1), Inches(3.1), Inches(2.0), Inches(0.3), f"Weight: {weight}", 11, color, True)

add_text_box(slide, Inches(0.8), Inches(3.9), Inches(6), Inches(0.4),
             "Qualitative Assessment (Days 13-15)", 18, ACCENT_PURPLE, True)

qual = [
    ("Judgment Test", "Present 3 clients from your Pond. Rank by potential value & justify.", ACCENT_BLUE),
    ("Coaching Simulation", "Role-play: Coach a new agent on handling an angry VIP client.", ACCENT_GREEN),
    ("Scenario Response", "High-value client hasn't responded in 5 days. What's your action plan?", ACCENT_PURPLE),
]

for i, (title, desc, color) in enumerate(qual):
    left = Inches(0.8 + i * 4.1)
    add_shape_with_fill(slide, left, Inches(4.4), Inches(3.8), Inches(1.4), CARD_BG, CARD_BORDER)
    add_rect(slide, left, Inches(4.4), Inches(3.8), Pt(4), color)
    add_text_box(slide, left + Inches(0.15), Inches(4.5), Inches(3.5), Inches(0.3), title, 14, color, True)
    add_text_box(slide, left + Inches(0.15), Inches(4.85), Inches(3.5), Inches(0.8), desc, 12, LIGHT_GRAY)

add_shape_with_fill(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8), ACCENT_GREEN)
add_text_box(slide, Inches(1.0), Inches(6.3), Inches(11), Inches(0.6),
             "PROMOTION THRESHOLD: Score >= 85/100  |  Top 4-6 scorers above threshold are promoted to Team Leader",
             18, WHITE, True, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 9: Risks & Resources
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "RESOURCES & RISKS", 12, ACCENT_YELLOW, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
             "Critical Success Factors & Risk Mitigation", 32, WHITE, True)

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(6), Inches(0.4),
             "Critical Success Factors", 18, ACCENT_GREEN, True)

factors = [
    ("120+ Leads/Day", "Qualified lead pipeline minimum", ACCENT_BLUE),
    ("2,000 Cold Pool", "Pre-qualified dormant clients", ACCENT_GREEN),
    ("1 Mgr / 10 Agents", "100% dedicated for 15 days", ACCENT_PURPLE),
    ("Tech Stack Ready", "CRM + auto-freeze + dashboards", ACCENT_YELLOW),
]

for i, (title, desc, color) in enumerate(factors):
    left = Inches(0.8 + i * 3.1)
    add_shape_with_fill(slide, left, Inches(2.1), Inches(2.8), Inches(1.0), CARD_BG, CARD_BORDER)
    add_rect(slide, left, Inches(2.1), Inches(2.8), Pt(4), color)
    add_text_box(slide, left + Inches(0.15), Inches(2.2), Inches(2.5), Inches(0.35), title, 16, WHITE, True)
    add_text_box(slide, left + Inches(0.15), Inches(2.6), Inches(2.5), Inches(0.4), desc, 11, LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(3.5), Inches(6), Inches(0.4),
             "Risk Mitigation", 18, ACCENT_RED, True)

risks = [
    ("Lead Shortage", "MEDIUM", "CRITICAL", "Secure 2-week buffer before start. 2 backup sources.", ACCENT_YELLOW),
    ("Agent Burnout", "HIGH", "HIGH", "Daily 1:1 check-ins. Small daily rewards. Clear promotion path.", ACCENT_RED),
    ("Quality Erosion", "MEDIUM", "HIGH", "Daily QA with same-day feedback. Peer listening sessions.", ACCENT_YELLOW),
    ("System Failure", "LOW", "CRITICAL", "Manual tracking sheets ready. 1 IT support on-call.", ACCENT_BLUE),
]

for i, (risk, prob, impact, mitigation, color) in enumerate(risks):
    top = Inches(4.0 + i * 0.85)
    add_shape_with_fill(slide, Inches(0.8), top, Inches(11.7), Inches(0.75), CARD_BG, CARD_BORDER)
    add_rect(slide, Inches(0.8), top, Pt(4), Inches(0.75), color)
    add_text_box(slide, Inches(1.1), top + Pt(5), Inches(2), Inches(0.3), risk, 14, WHITE, True)
    add_text_box(slide, Inches(3.2), top + Pt(5), Inches(1.5), Inches(0.3), f"Prob: {prob}", 11, color)
    add_text_box(slide, Inches(4.7), top + Pt(5), Inches(1.5), Inches(0.3), f"Impact: {impact}", 11, ACCENT_RED)
    add_text_box(slide, Inches(6.3), top + Pt(5), Inches(6), Inches(0.5), mitigation, 11, LIGHT_GRAY)


# ============================================================
# SLIDE 10: Scalability
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "POST-SPRINT SCALE", 12, ACCENT_PURPLE, True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
             "Scalability Model: 10 to 90 in 45 Days", 32, WHITE, True)

stages = [
    ("Day 1-15", "10", "Agents", "Sprint #1: Build the first cohort", ACCENT_BLUE),
    ("Day 16-30", "30", "Agents", "4-6 new TLs each manage 3-6 new agents", ACCENT_GREEN),
    ("Day 31-45", "90", "Agents", "Second wave TLs repeat the cycle", ACCENT_PURPLE),
]

for i, (period, count, unit, desc, color) in enumerate(stages):
    left = Inches(0.8 + i * 4.2)
    add_shape_with_fill(slide, left, Inches(1.8), Inches(3.8), Inches(2.4), CARD_BG, CARD_BORDER)
    add_shape_with_fill(slide, left, Inches(1.8), Inches(3.8), Inches(0.6), color)
    add_text_box(slide, left + Inches(0.2), Inches(1.88), Inches(3.4), Inches(0.4),
                 period, 18, WHITE, True)
    add_text_box(slide, left + Inches(0.2), Inches(2.55), Inches(3.4), Inches(0.6),
                 count, 48, WHITE, True)
    add_text_box(slide, left + Inches(0.2), Inches(3.2), Inches(3.4), Inches(0.3),
                 unit, 16, color, True)
    add_text_box(slide, left + Inches(0.2), Inches(3.55), Inches(3.4), Inches(0.5),
                 desc, 12, LIGHT_GRAY)
    if i < 2:
        add_text_box(slide, left + Inches(4.0), Inches(2.6), Inches(0.3), Inches(0.5),
                     ">>", 24, color, True)

add_text_box(slide, Inches(0.8), Inches(4.6), Inches(6), Inches(0.4),
             "Established Team Steady-State Metrics", 18, WHITE, True)

steady = [
    ("Pond per Agent", "80-120", "Steady state", ACCENT_BLUE),
    ("Daily Net Growth", "+2 to +3", "Maintenance mode", ACCENT_GREEN),
    ("Reactivation Rate", ">=15%", "Dormant recovery", ACCENT_PURPLE),
    ("TL Span of Control", "6-8", "Agents per TL", ACCENT_YELLOW),
]

for i, (title, value, desc, color) in enumerate(steady):
    left = Inches(0.8 + i * 3.1)
    add_metric_card(slide, left, Inches(5.1), Inches(2.8), Inches(1.4), title, value, desc, color)


# Save
output_path = r"C:\Users\us\Desktop\15_Day_Sprint_TL_Brief.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
